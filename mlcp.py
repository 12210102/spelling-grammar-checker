import streamlit as st
import requests
from docx import Document
from pptx import Presentation
import pandas as pd
import PyPDF2
import pdfplumber
from docx.opc.exceptions import PackageNotFoundError

API_KEY = 'cac6f7e7emsh46cbf311ae5026fp11c1e5jsne90bfffe277c'

# Function to check text via API
def check_text(text):
    headers = {
        'Authorization': f'Bearer {API_KEY}',
        'Content-Type': 'application/json'
    }
    data = {
        'text': text,
    }
    try:
       response = {"suggestions": [{"incorrect": "nmber", "correct": "number",}]}  # Example response
        
       return response
    
    except requests.exceptions.HTTPError as http_err:
        st.error(f"HTTP error occurred: {http_err}")
    except Exception as err:
        st.error(f"Other error occurred: {err}")
    return None

# Extract text from .pptx (PowerPoint) files
def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Extract text from .xlsx (Excel) files
def extract_text_from_excel(file):
    df = pd.read_excel(file, engine='openpyxl')  # Use openpyxl to read .xlsx files
    text = "\n".join(df.astype(str).values.flatten())  # Convert to string and flatten data
    return text

# Extract text from .pdf files using PyPDF2
def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# Alternative: Extract text from PDF using pdfplumber (for better accuracy)
def extract_text_from_pdf_plumber(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text

# Function to correct file based on file type
def correct_file(file, file_type):
    if file_type == "txt":
        text = file.read().decode("utf-8")
    elif file_type == "docx":
        try:
            doc = Document(file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except PackageNotFoundError:
            st.error("The uploaded file is not a valid .docx file or is corrupted.")
            return None
    elif file_type == "pptx":
        try:
            text = extract_text_from_pptx(file)
        except Exception as e:
            st.error("Error processing the .pptx file: " + str(e))
            return None
    elif file_type == "xlsx":
        try:
            text = extract_text_from_excel(file)
        except Exception as e:
            st.error("Error processing the .xlsx file: " + str(e))
            return None
    elif file_type == "pdf":
        try:
            text = extract_text_from_pdf(file)  # You can switch to pdfplumber if preferred
        except Exception as e:
            st.error("Error processing the .pdf file: " + str(e))
            return None
    else:
        st.error(f"File type '{file_type}' is not supported.")
        return None
    return check_text(text)

# Streamlit UI for Spelling and Grammar Check App
st.title("Spelling and Grammar Check App")

# Text input
text_input = st.text_area("Enter text to check:")
results = None
if st.button("Check text"):
    if text_input:
        results = check_text(text_input)
        st.write(results)

# File uploader that accepts multiple file types
uploaded_file = st.file_uploader("Upload a text, docx, pptx, xlsx, or pdf file", type=["txt", "docx", "pptx", "xlsx", "pdf"])
if uploaded_file is not None:
    file_type = uploaded_file.name.split('.')[-1]
    corrections = correct_file(uploaded_file, file_type)
    if corrections is not None:
        st.write(corrections)

# Suggestions and corrections
if results and 'suggestions' in results:
    selected_suggestion = st.selectbox(
        "Select a suggestion to correct:",
        [f"{s['incorrect']} -> {s['correct']}" for s in results['suggestions']]
    )
    
    if st.button("Apply Suggestion"):
        for suggestion in results['suggestions']:
            text_input = text_input.replace(suggestion['incorrect'], suggestion['correct'])
        st.text_area("Corrected Text:", value=text_input, height=300)
